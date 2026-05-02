"""
text_extractor.py
-----------------
Extracts text from all supported file types:
  - TXT / CSV  : plain read
  - PDF        : PyPDF2
  - DOCX       : python-docx
  - XLSX       : openpyxl
  - PPTX       : python-pptx
  - EML        : Python stdlib email module
  - MSG        : extract-msg library
  - ZIP        : extracts and reads text files inside the archive
  - Images     : pytesseract OCR (optional)
"""

import logging
import os

logger = logging.getLogger(__name__)

# ── Tesseract path (Windows) ─────────────────────────────────────────────────
_TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
if os.path.exists(_TESSERACT_PATH):
    try:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_PATH
    except ImportError:
        pass

SUPPORTED_EXTENSIONS = {
    ".txt", ".csv", ".pdf", ".docx",
    ".xlsx", ".pptx",
    ".eml", ".msg", ".zip",
    ".png", ".jpg", ".jpeg",
}


def extract_from_txt(filepath: str) -> str:
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.error("TXT extraction failed for %s: %s", filepath, e)
        return ""


def extract_from_pdf(filepath: str) -> str:
    try:
        import PyPDF2
        parts = []
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        logger.error("PDF extraction failed for %s: %s", filepath, e)
        return ""


def extract_from_docx(filepath: str) -> str:
    try:
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.error("DOCX extraction failed for %s: %s", filepath, e)
        return ""


def extract_from_xlsx(filepath: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
        parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        parts.append(str(cell))
        wb.close()
        return " ".join(parts)
    except Exception as e:
        logger.error("XLSX extraction failed for %s: %s", filepath, e)
        return ""


def extract_from_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        logger.error("PPTX extraction failed for %s: %s", filepath, e)
        return ""


def extract_from_eml(filepath: str) -> str:
    """Extract subject + body from a .eml email file using stdlib."""
    try:
        import email
        from email import policy
        with open(filepath, "rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)

        parts = []
        # Subject and sender give good classification signals
        if msg.get("subject"):
            parts.append(msg["subject"])
        if msg.get("from"):
            parts.append(msg["from"])

        # Extract plain-text body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        parts.append(payload.decode("utf-8", errors="ignore"))
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                parts.append(payload.decode("utf-8", errors="ignore"))

        return "\n".join(parts)
    except Exception as e:
        logger.error("EML extraction failed for %s: %s", filepath, e)
        return ""


def extract_from_msg(filepath: str) -> str:
    """Extract subject + body from an Outlook .msg file using extract-msg."""
    try:
        import extract_msg
        with extract_msg.openMsg(filepath) as msg:
            parts = []
            if msg.subject:
                parts.append(msg.subject)
            if msg.sender:
                parts.append(msg.sender)
            if msg.body:
                parts.append(msg.body)
            return "\n".join(str(p) for p in parts if p)
    except ImportError:
        logger.warning(
            "extract-msg not installed — skipping MSG for %s. "
            "Install with: pip install extract-msg", filepath
        )
        return ""
    except Exception as e:
        logger.error("MSG extraction failed for %s: %s", filepath, e)
        return ""


def extract_from_zip(filepath: str) -> str:
    """
    Open a ZIP archive and extract text from any supported text files inside.
    Reads up to 10 text files to keep processing fast.
    """
    try:
        import zipfile
        parts = []
        text_exts = {".txt", ".csv", ".md", ".py", ".js", ".html", ".xml", ".json"}
        count = 0

        with zipfile.ZipFile(filepath, "r") as zf:
            for name in zf.namelist():
                if count >= 10:
                    break
                ext = "." + name.lower().rsplit(".", 1)[-1] if "." in name else ""
                if ext in text_exts:
                    try:
                        with zf.open(name) as inner:
                            content = inner.read().decode("utf-8", errors="ignore")
                            parts.append(content[:2000])  # cap per file
                            count += 1
                    except Exception:
                        continue

        return "\n".join(parts)
    except Exception as e:
        logger.error("ZIP extraction failed for %s: %s", filepath, e)
        return ""


def extract_from_image(filepath: str) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(filepath)
        return pytesseract.image_to_string(img)
    except ImportError:
        logger.warning(
            "pytesseract / Pillow not installed — skipping OCR for %s.", filepath
        )
        return ""
    except Exception as e:
        logger.error("Image OCR failed for %s: %s", filepath, e)
        return ""


def extract_text(filepath: str) -> str:
    """Route to the correct extractor based on file extension."""
    ext = "." + filepath.lower().rsplit(".", 1)[-1] if "." in filepath else ""

    dispatch = {
        ".txt":  extract_from_txt,
        ".csv":  extract_from_txt,
        ".pdf":  extract_from_pdf,
        ".docx": extract_from_docx,
        ".xlsx": extract_from_xlsx,
        ".pptx": extract_from_pptx,
        ".eml":  extract_from_eml,
        ".msg":  extract_from_msg,
        ".zip":  extract_from_zip,
        ".png":  extract_from_image,
        ".jpg":  extract_from_image,
        ".jpeg": extract_from_image,
    }

    handler = dispatch.get(ext)
    if handler is None:
        logger.warning("Unsupported file type: %s", filepath)
        return ""
    return handler(filepath)
