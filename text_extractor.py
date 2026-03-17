"""
text_extractor.py
-----------------
Extracts text from supported file types:
  - TXT  : plain read
  - PDF  : PyPDF2
  - DOCX : python-docx
  - XLSX : openpyxl  (reads all cell values)
  - PPTX : python-pptx (reads all slide text frames)
  - CSV  : plain read (same as TXT)
  - Images (.png/.jpg/.jpeg) : pytesseract OCR (optional — skipped gracefully if not installed)

Each function returns extracted text as a plain string.
Empty string is returned on any failure so the pipeline can continue.
"""

import logging
import os

logger = logging.getLogger(__name__)

# ── Tesseract path (Windows) ─────────────────────────────────────────────────
# Automatically set the Tesseract path on Windows if it exists
_TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
if os.path.exists(_TESSERACT_PATH):
    try:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_PATH
    except ImportError:
        pass

SUPPORTED_EXTENSIONS = {
    ".txt", ".pdf", ".docx",
    ".xlsx", ".pptx", ".csv",
    ".png", ".jpg", ".jpeg",
}


# ── plain text / CSV ────────────────────────────────────────────────────────
def extract_from_txt(filepath: str) -> str:
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.error("TXT extraction failed for %s: %s", filepath, e)
        return ""


# ── PDF ─────────────────────────────────────────────────────────────────────
def extract_from_pdf(filepath: str) -> str:
    try:
        import PyPDF2
        parts = []
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        logger.error("PDF extraction failed for %s: %s", filepath, e)
        return ""


# ── DOCX ────────────────────────────────────────────────────────────────────
def extract_from_docx(filepath: str) -> str:
    try:
        from docx import Document
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.error("DOCX extraction failed for %s: %s", filepath, e)
        return ""


# ── XLSX ────────────────────────────────────────────────────────────────────
def extract_from_xlsx(filepath: str) -> str:
    """Read all non-empty cell values from every sheet."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
        parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        parts.append(str(cell))
        wb.close()
        return " ".join(parts)
    except Exception as e:
        logger.error("XLSX extraction failed for %s: %s", filepath, e)
        return ""


# ── PPTX ────────────────────────────────────────────────────────────────────
def extract_from_pptx(filepath: str) -> str:
    """Read all text frames from every slide."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        logger.error("PPTX extraction failed for %s: %s", filepath, e)
        return ""


# ── Images (OCR) ────────────────────────────────────────────────────────────
def extract_from_image(filepath: str) -> str:
    """
    Use pytesseract OCR to extract text from an image.
    If pytesseract or Tesseract is not installed, returns empty string
    and logs a warning (does not crash the pipeline).
    """
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(filepath)
        return pytesseract.image_to_string(img)
    except ImportError:
        logger.warning(
            "pytesseract / Pillow not installed — skipping OCR for %s. "
            "Install with: pip install pytesseract pillow  and install Tesseract.",
            filepath,
        )
        return ""
    except Exception as e:
        logger.error("Image OCR failed for %s: %s", filepath, e)
        return ""


# ── dispatcher ──────────────────────────────────────────────────────────────
def extract_text(filepath: str) -> str:
    """
    Route to the correct extractor based on file extension.

    Returns
    -------
    Extracted text string, or empty string on failure / unsupported type.
    """
    ext = filepath.lower().rsplit(".", 1)[-1]
    ext = f".{ext}"

    dispatch = {
        ".txt":  extract_from_txt,
        ".csv":  extract_from_txt,     # CSV is plain text
        ".pdf":  extract_from_pdf,
        ".docx": extract_from_docx,
        ".xlsx": extract_from_xlsx,
        ".pptx": extract_from_pptx,
        ".png":  extract_from_image,
        ".jpg":  extract_from_image,
        ".jpeg": extract_from_image,
    }

    handler = dispatch.get(ext)
    if handler is None:
        logger.warning("Unsupported file type: %s", filepath)
        return ""

    return handler(filepath)
